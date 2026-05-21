"""Build the 10-slide slideument PPTX for MLE Assignment 1."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------- palette (sampled from original deck) ----------------
CREAM = RGBColor(0xF1, 0xEF, 0xED)
TEAL = RGBColor(0x54, 0xBB, 0xC6)
SAGE = RGBColor(0xBF, 0xD1, 0x92)
LAV = RGBColor(0xEE, 0xDE, 0xEE)
PERI = RGBColor(0xBA, 0xCD, 0xE5)
DARK = RGBColor(0x1A, 0x2A, 0x3A)
TXT = RGBColor(0x2A, 0x2A, 0x2A)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
BRONZE = RGBColor(0xD9, 0xB3, 0x82)
SILVER_C = RGBColor(0xC8, 0xCF, 0xD6)
GOLD = RGBColor(0xE8, 0xC9, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9 slide @ 13.333" x 7.5"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

BLANK_LAYOUT = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = CREAM
    bg.shadow.inherit = False
    return s


def add_textbox(slide, left, top, width, height, text, *, size=14, bold=False,
                color=TXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        lines = [text]
    else:
        lines = text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *, size=14, color=TXT,
                bullet_color=DARK, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def add_rect(slide, left, top, width, height, fill, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_header(slide, slide_no, title, subtitle=None):
    # decorative accent bar
    add_rect(slide, Inches(0.5), Inches(0.55), Inches(0.18), Inches(0.55), TEAL)
    # slide number pill
    pill = add_rect(slide, Inches(0.78), Inches(0.55), Inches(0.55), Inches(0.4), SAGE, radius=True)
    add_textbox(
        slide, Inches(0.78), Inches(0.52), Inches(0.55), Inches(0.45),
        f"{slide_no:02d}", size=14, bold=True, color=DARK,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide, Inches(1.45), Inches(0.45), Inches(10), Inches(0.55),
        title, size=26, bold=True, color=DARK,
    )
    if subtitle:
        add_textbox(
            slide, Inches(1.45), Inches(1.02), Inches(10), Inches(0.35),
            subtitle, size=12, color=MUTED,
        )


def add_footer(slide, slide_no):
    add_textbox(
        slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
        "CS611 · MLE Assignment 1 · Medallion Pipeline · Jeslyn Wangsa",
        size=9, color=MUTED,
    )
    add_textbox(
        slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3),
        f"{slide_no} / 10", size=9, color=MUTED, align=PP_ALIGN.RIGHT,
    )


def add_table(slide, left, top, width, height, rows, *, header_fill=TEAL,
              header_text=WHITE, alt_fill=LAV, body_text=TXT,
              col_widths=None, font_size=11, header_size=12):
    n_rows = len(rows)
    n_cols = len(rows[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = int(width * w / total)
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else alt_fill
            tf = cell.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = ""
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Calibri"
            run.font.size = Pt(header_size if r_idx == 0 else font_size)
            run.font.bold = (r_idx == 0)
            run.font.color.rgb = header_text if r_idx == 0 else body_text
    return tbl_shape


def add_blob(slide, left, top, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # soften with transparency via XML
    sp = shape.fill.fore_color._xFill
    alpha = etree.SubElement(sp.find(qn('a:srgbClr')), qn('a:alpha'))
    alpha.set('val', '35000')
    return shape


# ==================== SLIDE 1 — Title ====================
s = add_slide()
# Decorative pastel blobs (top-right corner)
add_blob(s, Inches(9.0), Inches(-1.0), Inches(5.5), Inches(5.5), TEAL)
add_blob(s, Inches(10.5), Inches(2.5), Inches(4.0), Inches(4.0), SAGE)
add_blob(s, Inches(8.0), Inches(4.0), Inches(3.5), Inches(3.5), LAV)

add_textbox(s, Inches(0.7), Inches(0.5), Inches(6), Inches(0.5),
            "CS611 · Machine Learning Engineering", size=12, color=MUTED, bold=True)

add_textbox(s, Inches(0.7), Inches(1.8), Inches(11), Inches(1.2),
            "Loan Default", size=72, bold=True, color=DARK)
add_textbox(s, Inches(0.7), Inches(3.0), Inches(11), Inches(0.9),
            "Medallion Architecture Strategy & Implementation",
            size=28, color=TEAL, bold=True)

# TL;DR card
add_rect(s, Inches(0.7), Inches(4.5), Inches(9.0), Inches(1.4), WHITE, line=DARK, radius=True)
add_textbox(s, Inches(0.95), Inches(4.6), Inches(8.5), Inches(0.4),
            "TL;DR", size=12, bold=True, color=TEAL)
add_textbox(s, Inches(0.95), Inches(4.95), Inches(8.5), Inches(0.95),
            "A monthly Bronze→Silver→Gold pipeline that turns 4 raw CSVs into an ML-ready "
            "feature store and label store for predicting loan default at application time.",
            size=15, color=TXT)

add_textbox(s, Inches(0.7), Inches(6.4), Inches(8), Inches(0.4),
            "By Jeslyn Wangsa  ·  May 2026", size=12, color=MUTED, bold=True)
add_footer(s, 1)

# ==================== SLIDE 2 — Problem & success ====================
s = add_slide()
add_header(s, 2, "The problem we're solving",
           "What we need, what 'done' looks like, who reads this deck")

# Two columns
add_rect(s, Inches(0.5), Inches(1.6), Inches(6.0), Inches(4.9), WHITE, radius=True)
add_textbox(s, Inches(0.75), Inches(1.75), Inches(5.5), Inches(0.4),
            "Business problem", size=16, bold=True, color=TEAL)
add_bullets(s, Inches(0.75), Inches(2.25), Inches(5.5), Inches(4.0), [
    "Bank lends cash loans; some borrowers default.",
    "We need to score default risk at the moment of loan application, before any repayment data exists.",
    "An ML model will do the scoring (next assignment); this project ships the data pipeline that feeds it.",
    "Audience: ML engineers (will train the model), data engineers (will operate the pipeline), risk & business stakeholders (need to trust the inputs).",
], size=12)

add_rect(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(4.9), WHITE, radius=True)
add_textbox(s, Inches(7.05), Inches(1.75), Inches(5.5), Inches(0.4),
            "Definition of done", size=16, bold=True, color=TEAL)
add_bullets(s, Inches(7.05), Inches(2.25), Inches(5.5), Inches(4.0), [
    "Three medallion layers — Bronze (raw), Silver (cleaned), Gold (ML-ready) — for every source.",
    "Monthly snapshots from Jan 2023 to Dec 2024, idempotent and re-runnable.",
    "One Gold feature store + one Gold label store, joinable on (Customer_ID, snapshot_date).",
    "Zero target leakage: no post-application or repayment fields leak into features.",
    "Reproducible via docker-compose up and python main.py.",
], size=12)

add_footer(s, 2)

# ==================== SLIDE 3 — Data sources ====================
s = add_slide()
add_header(s, 3, "Raw data sources",
           "Four CSVs, semicolon-delimited, monthly snapshots Jan 2023 – Dec 2024")

rows = [
    ["File", "Role", "Grain", "Key columns"],
    ["feature_clickstream.csv", "Feature", "Customer × month",
     "Customer_ID, snapshot_date, fe_1 … fe_20"],
    ["feature_attributes.csv", "Feature", "Customer × month",
     "Customer_ID, Name, Age, SSN, Occupation, snapshot_date"],
    ["feature_financials.csv", "Feature", "Customer × month",
     "Customer_ID, Annual_Income, Credit_Mix, Outstanding_Debt, EMI, snapshot_date"],
    ["lms_loan_daily.csv", "Label only", "Loan × day",
     "loan_id, Customer_ID, installment_num, overdue_amt, due_amt, snapshot_date"],
]
add_table(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(3.6), rows,
          col_widths=[2.6, 1.2, 2.2, 6.3], font_size=12, header_size=13)

# Callout box
add_rect(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.2), LAV, radius=True)
add_textbox(s, Inches(0.75), Inches(5.7), Inches(11.8), Inches(0.4),
            "Why these splits matter for design", size=13, bold=True, color=DARK)
add_textbox(s, Inches(0.75), Inches(6.1), Inches(11.8), Inches(0.7),
            "Three feature files have very different update cadences and quality profiles, so each gets its own "
            "Bronze / Silver pipeline. The loan file is repayment history — usable only for labels, never for features.",
            size=11, color=TXT)
add_footer(s, 3)

# ==================== SLIDE 4 — Architecture diagram ====================
s = add_slide()
add_header(s, 4, "End-to-end architecture",
           "Raw drop-zone → immutable Bronze → cleaned Silver → ML-ready Gold → downstream ML")

s.shapes.add_picture(
    r"D:\SMU\Courses\Machine Learning Engineer\MLE_ASM1\deck\architecture.png",
    Inches(1.17), Inches(1.45), width=Inches(11.0),
)

# Bottom legend strip
y = Inches(6.35)
items = [
    ("CSV", BRONZE, "Raw + Bronze: human-readable, easy to inspect"),
    ("Parquet", SILVER_C, "Silver + Gold: typed, compressed, columnar"),
    ("Monthly partition", GOLD, "One file per snapshot_date — incremental runs"),
]
xpos = 0.5
for tag, color, desc in items:
    add_rect(s, Inches(xpos), y, Inches(0.18), Inches(0.4), color)
    add_textbox(s, Inches(xpos + 0.3), Inches(6.32), Inches(1.7), Inches(0.25),
                tag, size=11, bold=True, color=DARK)
    add_textbox(s, Inches(xpos + 0.3), Inches(6.55), Inches(4), Inches(0.3),
                desc, size=9, color=MUTED)
    xpos += 4.2
add_footer(s, 4)

# ==================== SLIDE 5 — Bronze ====================
s = add_slide()
add_header(s, 5, "Bronze — raw landing zone",
           "Slice once, store immutably, never transform")

# Left: design choices
add_textbox(s, Inches(0.5), Inches(1.7), Inches(6.5), Inches(0.4),
            "Three design choices", size=16, bold=True, color=TEAL)
add_bullets(s, Inches(0.5), Inches(2.15), Inches(6.5), Inches(4.5), [
    "Monthly partitions — one file per snapshot_date. Reprocessing is incremental, not full-history.",
    "CSV format preserved — keeps Bronze byte-identical to the source system so Silver fixes can never poison the raw record.",
    "Immutable on write — Bronze is the audit trail. Any Silver or Gold version can be reproduced from a frozen Bronze partition.",
], size=12, line_spacing=1.4)

# Right: filename pattern card
add_rect(s, Inches(7.3), Inches(1.7), Inches(5.5), Inches(4.5), WHITE, radius=True, line=SAGE)
add_textbox(s, Inches(7.55), Inches(1.85), Inches(5), Inches(0.4),
            "Filename pattern", size=14, bold=True, color=TEAL)
add_textbox(s, Inches(7.55), Inches(2.3), Inches(5), Inches(0.4),
            "datamart/bronze/<source>/", size=11, color=TXT, font="Consolas")
add_textbox(s, Inches(7.55), Inches(2.6), Inches(5), Inches(0.4),
            "  bronze_<source>_YYYY_MM_01.csv", size=11, color=TXT, font="Consolas")

add_textbox(s, Inches(7.55), Inches(3.2), Inches(5), Inches(0.4),
            "Example", size=12, bold=True, color=DARK)
add_textbox(s, Inches(7.55), Inches(3.55), Inches(5), Inches(0.4),
            "datamart/bronze/clickstream/", size=10.5, color=MUTED, font="Consolas")
add_textbox(s, Inches(7.55), Inches(3.85), Inches(5), Inches(0.4),
            "  bronze_clickstream_2023_01_01.csv", size=10.5, color=MUTED, font="Consolas")

add_textbox(s, Inches(7.55), Inches(4.5), Inches(5), Inches(0.4),
            "Driven by", size=12, bold=True, color=DARK)
add_textbox(s, Inches(7.55), Inches(4.85), Inches(5), Inches(0.4),
            "utils/data_processing_bronze_table.py", size=10.5, color=MUTED, font="Consolas")
add_textbox(s, Inches(7.55), Inches(5.2), Inches(5), Inches(0.6),
            "main.py loops Jan-2023 → Dec-2024 and writes one partition per source per month.",
            size=10.5, color=TXT)

# Bottom note
add_rect(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5), SAGE, radius=True)
add_textbox(s, Inches(0.75), Inches(6.43), Inches(12), Inches(0.45),
            "No type casts, no filters, no joins at Bronze. The only operation is 'select * where snapshot_date = X'.",
            size=11, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 5)

# ==================== SLIDE 6 — Silver ====================
s = add_slide()
add_header(s, 6, "Silver — typed, cleaned, trustworthy",
           "Parquet partitions with enforced schema and explicit data-quality rules")

# Table of cleaning rules
rows = [
    ["Issue found in raw", "Treatment", "Why"],
    ["Mixed / wrong data types", "Cast each column to its target type (Int, Float, Date, String)",
     "Downstream feature engineering needs a stable schema"],
    ["PII — SSN, Name", "Drop SSN; keep Name only if needed for ops",
     "Legal & privacy risk; SSN has no predictive value"],
    ["Corrupted strings (\"_______\", \"!@9#%8\")", "Drop the entire row",
     "Looks like upstream encoding failure, not a missing value"],
    ["Trailing underscores (\"45_\", \"_\")", "Strip and re-cast",
     "Recoverable noise — preserve the signal"],
    ["Numeric outliers (Age = 8000, Income < 0)", "Tukey IQR fences + domain bounds (Age 18–70)",
     "Realistic loan applicants only; protects model from leverage points"],
    ["Loan repayment fields (lms)", "Compute mob (month-on-book) and dpd (days past due)",
     "Engineered once at Silver so Gold label logic stays trivial"],
]
add_table(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(4.6), rows,
          col_widths=[3.0, 4.0, 5.3], font_size=11, header_size=12)

add_rect(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5), LAV, radius=True)
add_textbox(s, Inches(0.75), Inches(6.43), Inches(12), Inches(0.45),
            "Output: datamart/silver/<source>/silver_<source>_YYYY_MM_01.parquet  ·  one directory per month, schema-on-write.",
            size=11, color=DARK, font="Calibri", bold=True, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 6)

# ==================== SLIDE 7 — Gold Feature Store ====================
s = add_slide()
add_header(s, 7, "Gold — feature store",
           "Two tables, deliberately split by update cadence")

# Left column: two table cards
def feat_card(left, top, title, lines, color):
    add_rect(s, left, top, Inches(6.0), Inches(2.4), WHITE, radius=True, line=color)
    add_rect(s, left, top, Inches(6.0), Inches(0.5), color)
    add_textbox(s, left + Inches(0.2), top + Inches(0.05), Inches(5.6), Inches(0.4),
                title, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    tb = add_textbox(s, left + Inches(0.2), top + Inches(0.6), Inches(5.6), Inches(1.8),
                     lines, size=11, color=TXT)


feat_card(Inches(0.5), Inches(1.7),
          "feature_store_clickstream.parquet",
          ["20 numeric behaviour features (fe_1 … fe_20)",
           "Grain: Customer_ID × snapshot_date",
           "Real-time-ish — refreshed independently from batch sources",
           "No engineered columns; trust the raw signal"],
          TEAL)

feat_card(Inches(6.8), Inches(1.7),
          "feature_store_customer.parquet",
          ["INNER JOIN: silver attributes ⟕ financials on (Customer_ID, snapshot_date)",
           "One-hot: Occupation, Credit_Mix, Payment_Behaviour",
           "Type_of_Loan exploded into n_loans_<type> counts",
           "Slow-changing snapshot — batch refresh on schedule"],
          SAGE)

# Why split section
add_rect(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(1.95), WHITE, radius=True, line=DARK)
add_textbox(s, Inches(0.75), Inches(4.4), Inches(12), Inches(0.4),
            "Why two tables, not one?", size=14, bold=True, color=TEAL)
add_bullets(s, Inches(0.75), Inches(4.85), Inches(12), Inches(1.4), [
    "Different update cadence — joining them at Gold would force the slower table to wait, or staleness to leak in.",
    "Different ownership — attributes/financials come from core banking; clickstream comes from web analytics.",
    "Joined on demand at training time via (Customer_ID, snapshot_date) — keeps storage normalised, joins explicit.",
], size=11, line_spacing=1.3)

# Schema preview strip
add_textbox(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.3),
            "Schema preview — feature_store_customer:  Customer_ID · snapshot_date · Age · Annual_Income · "
            "Outstanding_Debt · EMI · Credit_Mix_Good · Credit_Mix_Standard · n_loans_Personal · n_loans_Auto · …",
            size=10, color=MUTED, font="Consolas")
add_footer(s, 7)

# ==================== SLIDE 8 — Gold Label Store ====================
s = add_slide()
add_header(s, 8, "Gold — label store",
           "One row per loan, binary default flag observed 6 months after origination")

# Big rule card
add_rect(s, Inches(0.5), Inches(1.7), Inches(7.5), Inches(3.4), WHITE, radius=True, line=GOLD)
add_textbox(s, Inches(0.75), Inches(1.85), Inches(7), Inches(0.45),
            "Label definition", size=16, bold=True, color=TEAL)

add_textbox(s, Inches(0.75), Inches(2.4), Inches(7), Inches(0.4),
            "Observation window", size=12, bold=True, color=DARK)
add_textbox(s, Inches(0.75), Inches(2.75), Inches(7), Inches(0.4),
            "mob = 6   (the 6th payment after loan start)",
            size=13, color=TXT, font="Consolas")

add_textbox(s, Inches(0.75), Inches(3.3), Inches(7), Inches(0.4),
            "Default rule", size=12, bold=True, color=DARK)
add_textbox(s, Inches(0.75), Inches(3.65), Inches(7), Inches(0.4),
            "dpd  >= 30  →  label = 1   (defaulted)",
            size=13, color=TXT, font="Consolas")
add_textbox(s, Inches(0.75), Inches(4.0), Inches(7), Inches(0.4),
            "otherwise  →  label = 0   (current / paid)",
            size=13, color=TXT, font="Consolas")

add_textbox(s, Inches(0.75), Inches(4.55), Inches(7), Inches(0.4),
            "Columns retained", size=12, bold=True, color=DARK)
add_textbox(s, Inches(0.75), Inches(4.85), Inches(7), Inches(0.4),
            "loan_id · Customer_ID · label · label_def · snapshot_date",
            size=11.5, color=TXT, font="Consolas")

# Right: rationale + reuse note
add_rect(s, Inches(8.3), Inches(1.7), Inches(4.5), Inches(3.4), LAV, radius=True)
add_textbox(s, Inches(8.55), Inches(1.85), Inches(4.0), Inches(0.4),
            "Why mob = 6 ?", size=14, bold=True, color=DARK)
add_bullets(s, Inches(8.55), Inches(2.3), Inches(4.0), Inches(2.5), [
    "Short enough to label most loans in the 24-month window.",
    "Long enough that early-stage hiccups stop being noise.",
    "Standard industry vintage horizon for unsecured lending.",
], size=11, line_spacing=1.3)

# Reuse note
add_rect(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.9), SAGE, radius=True)
add_textbox(s, Inches(0.75), Inches(5.4), Inches(12), Inches(0.4),
            "Code reuse", size=13, bold=True, color=DARK)
add_textbox(s, Inches(0.75), Inches(5.75), Inches(12), Inches(0.4),
            "Label-store logic reused from Lab 2 verbatim — same Silver→Gold transformation on lms_loan_daily, "
            "same mob/dpd thresholds, same output schema.", size=11, color=DARK)

# Bottom: joinability note
add_textbox(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
            "Joinable to either feature_store table on (Customer_ID, snapshot_date) for ML training.",
            size=11, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 8)

# ==================== SLIDE 9 — Data leakage controls ====================
s = add_slide()
add_header(s, 9, "Data leakage controls",
           "Rule: every feature must be knowable strictly before the loan application date")

# Three columns: target / temporal / train-test
def leak_col(left, title, body):
    add_rect(s, left, Inches(1.7), Inches(4.1), Inches(3.5), WHITE, radius=True, line=TEAL)
    add_rect(s, left, Inches(1.7), Inches(4.1), Inches(0.5), TEAL)
    add_textbox(s, left + Inches(0.15), Inches(1.75), Inches(3.8), Inches(0.4),
                title, size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, left + Inches(0.15), Inches(2.3), Inches(3.8), Inches(2.8),
                body, size=10.5, line_spacing=1.25)


leak_col(Inches(0.5), "Target leakage", [
    "Drop entire lms_loan_daily from features — it IS the answer (repayment history).",
    "Drop Interest_Rate — set by the underwriter after seeing the risk score; encodes the label.",
    "Re-check any 'Num_of_Delayed_Payment'-style fields: kept only when measured pre-application.",
])

leak_col(Inches(4.65), "Temporal leakage", [
    "Snapshot_date is the as-of cut. We never join a feature row from t+1 to a loan applied at t.",
    "Monthly partitions enforce the cut physically: a Jan-2023 training run only reads ≤ 2023-01-01 files.",
    "Clickstream & financials are user-level snapshots — confirmed safe.",
])

leak_col(Inches(8.8), "Train-test contamination", [
    "Feature engineering is fit per partition, not across the whole table.",
    "No global means / scalers computed at Gold — that work belongs to the ML pipeline, after the split.",
    "Same Customer_ID can appear in train and test, but only at different snapshot_dates.",
])

# Bottom: sanity check note
add_rect(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5), LAV, radius=True)
add_textbox(s, Inches(0.75), Inches(5.5), Inches(12), Inches(0.4),
            "Sanity check (recommended by brief)", size=13, bold=True, color=DARK)
add_textbox(s, Inches(0.75), Inches(5.9), Inches(12), Inches(0.9),
            "Fit a simple logistic regression on Gold feature_store ⋈ label_store. If AUC is implausibly high "
            "(> 0.95) on held-out months, suspect a leaky column and trace it back through Silver.",
            size=11, color=TXT)
add_footer(s, 9)

# ==================== SLIDE 10 — Tech stack & how to run ====================
s = add_slide()
add_header(s, 10, "Tech stack, repo & how to run",
           "Reproducible end-to-end in one command")

# Tech stack chips
add_textbox(s, Inches(0.5), Inches(1.65), Inches(8), Inches(0.4),
            "Tech stack", size=14, bold=True, color=TEAL)
chips = [("Python 3", TEAL), ("PySpark", SAGE), ("Docker", LAV), ("JupyterLab", PERI),
         ("Parquet", GOLD), ("Git / GitHub", SILVER_C)]
xpos = 0.5
for label, color in chips:
    add_rect(s, Inches(xpos), Inches(2.15), Inches(1.55), Inches(0.45), color, radius=True)
    add_textbox(s, Inches(xpos), Inches(2.15), Inches(1.55), Inches(0.45),
                label, size=11, bold=True, color=DARK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    xpos += 1.7

# How to run card
add_rect(s, Inches(0.5), Inches(2.9), Inches(7.8), Inches(3.7), WHITE, radius=True, line=DARK)
add_textbox(s, Inches(0.75), Inches(3.0), Inches(7.3), Inches(0.4),
            "How to run", size=14, bold=True, color=TEAL)
steps = [
    ("1.  git clone <repo>  &&  cd <repo>", "clone the project"),
    ("2.  docker-compose build", "build the Spark + Jupyter image"),
    ("3.  docker-compose up", "launches JupyterLab on http://localhost:8888"),
    ("4.  python main.py", "loops Jan-23 → Dec-24, writes datamart/{bronze,silver,gold}"),
    ("5.  Inspect datamart/gold/ for ML-ready Parquet tables", "feature & label stores land here"),
]
y0 = 3.45
for cmd, note in steps:
    add_textbox(s, Inches(0.75), Inches(y0), Inches(7.3), Inches(0.3),
                cmd, size=11, bold=True, color=DARK, font="Consolas")
    add_textbox(s, Inches(0.75), Inches(y0 + 0.28), Inches(7.3), Inches(0.3),
                note, size=10, color=MUTED)
    y0 += 0.58

# Repo + author card
add_rect(s, Inches(8.6), Inches(2.9), Inches(4.3), Inches(3.7), LAV, radius=True)
add_textbox(s, Inches(8.85), Inches(3.0), Inches(3.8), Inches(0.4),
            "Repository", size=14, bold=True, color=DARK)
add_textbox(s, Inches(8.85), Inches(3.45), Inches(3.8), Inches(0.4),
            "github.com/<your-handle>/MLE_ASM1", size=11, color=TXT, font="Consolas")
add_textbox(s, Inches(8.85), Inches(3.8), Inches(3.8), Inches(0.4),
            "(see Readme.txt in submission zip)", size=9, color=MUTED)

add_textbox(s, Inches(8.85), Inches(4.5), Inches(3.8), Inches(0.4),
            "Deliverables in zip", size=13, bold=True, color=DARK)
add_bullets(s, Inches(8.85), Inches(4.9), Inches(3.8), Inches(1.6), [
    "main.py",
    "Dockerfile · docker-compose.yaml",
    "requirements.txt",
    "utils/ · data/",
    "Readme.txt (repo link)",
], size=10.5, line_spacing=1.15)

# Closing line
add_rect(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.25), TEAL)
add_footer(s, 10)


out = r"D:\SMU\Courses\Machine Learning Engineer\MLE_ASM1\deck\Assignment1_Deck.pptx"
prs.save(out)
print("saved", out)
