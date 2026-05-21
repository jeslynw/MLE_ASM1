"""Two standalone slides explaining the table split/merge reasoning."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CREAM = RGBColor(0xF1, 0xEF, 0xED)
TEAL = RGBColor(0x54, 0xBB, 0xC6)
SAGE = RGBColor(0xBF, 0xD1, 0x92)
LAV = RGBColor(0xEE, 0xDE, 0xEE)
PERI = RGBColor(0xBA, 0xCD, 0xE5)
DARK = RGBColor(0x1A, 0x2A, 0x3A)
TXT = RGBColor(0x2A, 0x2A, 0x2A)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xE8, 0xC9, 0x77)
RED = RGBColor(0xE0, 0x8B, 0x8B)
GREEN = RGBColor(0x7E, 0xB6, 0x7E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = CREAM
    bg.shadow.inherit = False
    return s


def tb(slide, left, top, w, h, text, *, size=12, bold=False, color=TXT,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    b = slide.shapes.add_textbox(left, top, w, h)
    f = b.text_frame
    for m in (f.margin_left, f.margin_right, f.margin_top, f.margin_bottom): pass
    f.margin_left = Emu(0); f.margin_right = Emu(0)
    f.margin_top = Emu(0); f.margin_bottom = Emu(0)
    f.word_wrap = True
    f.vertical_anchor = anchor
    lines = [text] if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = f.paragraphs[0] if i == 0 else f.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return b


def bullets(slide, left, top, w, h, items, *, size=11, color=TXT, line_spacing=1.25):
    b = slide.shapes.add_textbox(left, top, w, h)
    f = b.text_frame; f.word_wrap = True
    f.margin_left = Emu(0); f.margin_right = Emu(0)
    f.margin_top = Emu(0); f.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        p = f.paragraphs[0] if i == 0 else f.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = "•  " + it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return b


def rect(slide, left, top, w, h, fill, line=None, radius=False):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def header(slide, num, title, subtitle):
    rect(slide, Inches(0.5), Inches(0.55), Inches(0.18), Inches(0.55), TEAL)
    rect(slide, Inches(0.78), Inches(0.55), Inches(0.55), Inches(0.4), SAGE, radius=True)
    tb(slide, Inches(0.78), Inches(0.52), Inches(0.55), Inches(0.45),
       f"{num:02d}", size=14, bold=True, color=DARK,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(slide, Inches(1.45), Inches(0.45), Inches(11), Inches(0.55),
       title, size=26, bold=True, color=DARK)
    tb(slide, Inches(1.45), Inches(1.02), Inches(11), Inches(0.35),
       subtitle, size=12, color=MUTED)


def footer(slide, n):
    tb(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
       "CS611 · MLE Assignment 1 · Medallion Pipeline · Jeslyn Wangsa",
       size=9, color=MUTED)
    tb(slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3),
       f"{n} / 10", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def table(slide, left, top, w, h, rows, *, col_widths=None,
          header_fill=TEAL, header_text=WHITE, alt_fill=LAV,
          font_size=11, header_size=12):
    nr, nc = len(rows), len(rows[0])
    ts = slide.shapes.add_table(nr, nc, left, top, w, h)
    t = ts.table
    if col_widths:
        tot = sum(col_widths)
        for i, cw in enumerate(col_widths):
            t.columns[i].width = int(w * cw / tot)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else alt_fill
            tf = cell.text_frame; tf.word_wrap = True
            tf.paragraphs[0].text = ""
            p = tf.paragraphs[0]; r = p.add_run()
            r.text = str(val)
            r.font.name = "Calibri"
            r.font.size = Pt(header_size if ri == 0 else font_size)
            r.font.bold = (ri == 0)
            r.font.color.rgb = header_text if ri == 0 else TXT
            if ri > 0 and ci > 0 and ci < len(row):
                # colour-code the verdict cells if value is ✓ / ✗
                if val == "match":
                    r.font.color.rgb = GREEN; r.font.bold = True
                elif val == "differ":
                    r.font.color.rgb = RED; r.font.bold = True
    return ts


# =================== SLIDE A — Decision rule + compatibility matrix ===================
s = add_slide()
header(s, 7, "Gold table layout — the decision rule",
       "Why attributes + financials merge into one table, while clickstream stays apart")

# Decision rule banner
rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.85), DARK, radius=True)
tb(s, Inches(0.75), Inches(1.62), Inches(11.8), Inches(0.35),
   "Decision rule", size=11, bold=True, color=TEAL)
tb(s, Inches(0.75), Inches(1.92), Inches(11.8), Inches(0.45),
   "Tables stay separate when update cadence, ownership, grain, or reliability profile differ. "
   "They merge only when all four align AND they are nearly always consumed together.",
   size=13, color=WHITE, bold=True)

# Compatibility matrix
tb(s, Inches(0.5), Inches(2.6), Inches(8), Inches(0.4),
   "Compatibility matrix across the three feature sources",
   size=13, bold=True, color=DARK)
rows = [
    ["Dimension", "Attributes", "Financials", "Clickstream"],
    ["Update cadence", "Monthly batch", "Monthly batch", "Hourly stream"],
    ["System of record", "Core banking", "Core banking", "Web analytics"],
    ["Grain", "Customer × month", "Customer × month", "Customer × month"],
    ["Population coverage", "All customers", "All customers", "Web-active only"],
    ["Schema volatility", "Low", "Low", "High (new events)"],
    ["Verdict vs attributes", "—", "match", "differ"],
]
table(s, Inches(0.5), Inches(3.05), Inches(8.0), Inches(3.4), rows,
      col_widths=[2.2, 1.9, 1.9, 2.0], font_size=10.5, header_size=11.5)

# Right column: outcome
rect(s, Inches(8.8), Inches(3.05), Inches(4.0), Inches(3.4), WHITE, radius=True, line=TEAL)
tb(s, Inches(9.0), Inches(3.18), Inches(3.7), Inches(0.4),
   "Outcome", size=14, bold=True, color=TEAL)

# Two outcome cards inside
rect(s, Inches(9.0), Inches(3.65), Inches(3.6), Inches(1.25), SAGE, radius=True)
tb(s, Inches(9.15), Inches(3.7), Inches(3.4), Inches(0.35),
   "MERGE → feature_store_customer", size=11, bold=True, color=DARK)
tb(s, Inches(9.15), Inches(4.02), Inches(3.4), Inches(0.85),
   "attributes ⨝ financials. Four dimensions align; analysts always use them together.",
   size=10, color=DARK)

rect(s, Inches(9.0), Inches(5.05), Inches(3.6), Inches(1.3), PERI, radius=True)
tb(s, Inches(9.15), Inches(5.1), Inches(3.4), Inches(0.35),
   "SPLIT → feature_store_clickstream", size=11, bold=True, color=DARK)
tb(s, Inches(9.15), Inches(5.42), Inches(3.4), Inches(0.95),
   "Cadence, owner, and coverage all differ. Joining at Gold would inject staleness or silently drop rows.",
   size=10, color=DARK)

# Bottom: business hook
rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4), GOLD, radius=True)
tb(s, Inches(0.75), Inches(6.57), Inches(12), Inches(0.35),
   "Business framing: the split mirrors source-system boundaries — one table per (system of record × refresh cadence).",
   size=10.5, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 7)


# =================== SLIDE B — Why feature store ≠ label store ===================
s = add_slide()
header(s, 8, "Why feature store ≠ label store",
       "Decoupling features from outcomes protects against leakage and unblocks refresh cadence")

# Four reason cards in a 2×2 grid
def reason(left, top, title, body, color):
    rect(s, left, top, Inches(6.0), Inches(2.05), WHITE, radius=True, line=color)
    rect(s, left, top, Inches(0.18), Inches(2.05), color)
    tb(s, left + Inches(0.35), top + Inches(0.12), Inches(5.5), Inches(0.4),
       title, size=13, bold=True, color=DARK)
    tb(s, left + Inches(0.35), top + Inches(0.55), Inches(5.5), Inches(1.4),
       body, size=10.5, color=TXT)


reason(Inches(0.5), Inches(1.6),
       "Different time semantics",
       "Features are observed AT loan application time. Labels are observed 6 months LATER. "
       "Mixed in one table, any join on snapshot_date silently creates target leakage.",
       TEAL)

reason(Inches(6.83), Inches(1.6),
       "Different refresh triggers",
       "Features refresh as upstream snapshots land. Labels can't crystallise until mob=6 is reached. "
       "Coupling them means feature refresh would be blocked waiting for label maturity.",
       SAGE)

reason(Inches(0.5), Inches(3.8),
       "Different consumers",
       "The same feature store can feed multiple models — default at mob=3, churn, prepayment, "
       "fraud. Decoupling makes the feature store a reusable asset, not a one-model artefact.",
       LAV)

reason(Inches(6.83), Inches(3.8),
       "Different audit story",
       "Regulators ask: 'what did the model know at decision time?' That question has a clean "
       "answer only when features and outcomes live in separate tables with separate timestamps.",
       PERI)

# Bottom: explicit join illustration
rect(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.95), DARK, radius=True)
tb(s, Inches(0.75), Inches(6.13), Inches(11.8), Inches(0.35),
   "ML training time — leakage is forced into the open", size=11, bold=True, color=TEAL)
tb(s, Inches(0.75), Inches(6.45), Inches(11.8), Inches(0.5),
   "feature_store  ⨝  label_store  ON  Customer_ID, snapshot_date     "
   "→  an explicit, reviewable join. No leak can sneak in by accident.",
   size=11, color=WHITE, font="Consolas")

footer(s, 8)


out = r"D:\SMU\Courses\Machine Learning Engineer\MLE_ASM1\deck\Table_Design_Slides.pptx"
prs.save(out)
print("saved", out)
